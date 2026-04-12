
##  `create_presentation.py`

"""
Generate PowerPoint Presentation for EcoCrop Tunisia Defense
Run: python create_presentation.py
Output: PRESENTATION.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
GREEN_MED = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)
GREEN_PALE = RGBColor(0xE8, 0xF5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x61, 0x61, 0x61)
AMBER = RGBColor(0xFF, 0x8F, 0x00)
BLUE = RGBColor(0x15, 0x65, 0xC0)
RED = RGBColor(0xC6, 0x28, 0x28)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_DARK
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN_PALE if r % 2 == 0 else WHITE
    return table

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_shape_bg(slide1, Inches(0), Inches(0), Inches(13.333), Inches(7.5), GREEN_DARK)
add_shape_bg(slide1, Inches(0), Inches(5.5), Inches(13.333), Inches(0.08), GREEN_LIGHT)

add_text_box(slide1, Inches(1), Inches(0.5), Inches(11), Inches(0.6),
             "ECOCROP TUNISIA", 20, False, GREEN_LIGHT, PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
             "Crop Yield Prediction Using\nMachine Learning", 42, True, WHITE, PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
             "Predicting cereal production across 24 governorates from agro-climatic data (2016–2024)", 18, False, RGBColor(0xA5, 0xD6, 0xA7), PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1), Inches(4.5), Inches(5), Inches(0.5),
             "Ferdaws Saidi  &  Aya Gharsalli", 20, True, WHITE, PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1), Inches(5.8), Inches(5), Inches(0.5),
             "2024  |  Machine Learning Project", 14, False, RGBColor(0x81, 0xC7, 0x84), PP_ALIGN.LEFT)
add_text_box(slide1, Inches(8), Inches(5.8), Inches(4.5), Inches(0.8),
             "🌾", 60, False, WHITE, PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: Problem Statement & Objectives
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide2, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Problem Statement & Objectives", 30, True, WHITE, PP_ALIGN.LEFT)

add_text_box(slide2, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5),
             "🔥 The Problem", 22, True, GREEN_DARK)
add_bullet_list(slide2, Inches(0.5), Inches(2.1), Inches(5.8), Inches(3.5), [
    "• Tunisia's agriculture is climate-sensitive",
    "• Cereals = staple food for 11M+ people",
    "• Yield varies significantly by region & year",
    "• Traditional forecasting lacks precision",
    "• No ML-based decision support tool exists"
], 15, GRAY)

add_shape_bg(slide2, Inches(6.8), Inches(1.5), Inches(6), Inches(5.2), GREEN_PALE)
add_text_box(slide2, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
             "🎯 Our Objectives", 22, True, GREEN_DARK)
add_bullet_list(slide2, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0), [
    "1. Explore 9 years of agro-climatic data",
    "2. Engineer meaningful weather features",
    "3. Train & compare regression models",
    "4. Identify key yield drivers",
    "5. Deploy an interactive prediction tool",
    "6. Achieve R² > 0.90 on test set"
], 15, BLACK)

# ============================================================
# SLIDE 3: Dataset & Pipeline
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide3, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Dataset Overview & ML Pipeline", 30, True, WHITE, PP_ALIGN.LEFT)

# Dataset box
add_shape_bg(slide3, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.3), GREEN_PALE)
add_text_box(slide3, Inches(0.7), Inches(1.6), Inches(5), Inches(0.5),
             "📊 Dataset at a Glance", 20, True, GREEN_DARK)
data_items = [
    "• 1,479 observations × 15 columns",
    "• 24 governorates × 9 years (2016–2024)",
    "• 5 weather features: precipitation,",
    "  temperature, humidity, solar, wind",
    "• 7 crop yield targets (Tonnes)",
    "• Primary target: Cereales (T)",
    "• 0 missing values, 0 duplicates",
    "• Source: EcoCrop Tunisia (cleaned)"
]
add_bullet_list(slide3, Inches(0.7), Inches(2.2), Inches(5), Inches(4.0), data_items, 14, BLACK)

# Pipeline box
add_text_box(slide3, Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.5),
             "⚙️ ML Pipeline", 20, True, GREEN_DARK)
pipeline_steps = [
    ("1. Data Loading", "CSV import & validation"),
    ("2. EDA", "Distributions, correlations, outliers"),
    ("3. Feature Engineering", "Season, RainBin, TempZone, HeatIdx"),
    ("4. Encoding", "LabelEncoder for categoricals"),
    ("5. Scaling", "StandardScaler (z-score)"),
    ("6. Split", "80/20 train-test (seed=42)"),
    ("7. Modeling", "LR → DT → RF → RF-Tuned"),
    ("8. Evaluation", "R², RMSE, MAE, CV"),
    ("9. Analysis", "Feature importance, residuals"),
]
y_pos = 2.1
for step, desc in pipeline_steps:
    add_text_box(slide3, Inches(6.5), Inches(y_pos), Inches(2.5), Inches(0.4),
                 step, 13, True, GREEN_MED)
    add_text_box(slide3, Inches(9.0), Inches(y_pos), Inches(4), Inches(0.4),
                 desc, 13, False, GRAY)
    y_pos += 0.48

# ============================================================
# SLIDE 4: Model Comparison & Feature Importance
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide4, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Model Comparison & Feature Importance", 30, True, WHITE, PP_ALIGN.LEFT)

# Model comparison table
add_text_box(slide4, Inches(0.5), Inches(1.4), Inches(6), Inches(0.5),
             "📊 Model Performance", 20, True, GREEN_DARK)
model_data = [
    ["Model", "R² Test", "RMSE (T)", "MAE (T)", "Verdict"],
    ["Linear Regression", "0.25", "2,500", "1,800", "Too simple"],
    ["Decision Tree", "0.95", "550", "300", "Overfits"],
    ["Random Forest", "0.96", "480", "310", "Strong"],
    ["RF Tuned ⭐", "0.967", "451", "299", "Best"],
]
add_table(slide4, Inches(0.5), Inches(2.0), Inches(6.2), Inches(2.2),
          5, 5, model_data)

# Feature importance
add_text_box(slide4, Inches(7), Inches(1.4), Inches(6), Inches(0.5),
             "🌟 Top 5 Features", 20, True, GREEN_DARK)
feat_data = [
    ["Rank", "Feature", "Importance", "Why"],
    ["1", "Governorate", "28%", "Soil, altitude, microclimate"],
    ["2", "Heat Index", "18%", "Combined heat stress"],
    ["3", "Solar Radiation", "15%", "Photosynthesis driver"],
    ["4", "Air Temperature", "12%", "Direct thermal effect"],
    ["5", "Year", "8%", "Temporal trend"],
]
add_table(slide4, Inches(7), Inches(2.0), Inches(5.8), Inches(2.2),
          5, 4, feat_data)

# Key takeaway box
add_shape_bg(slide4, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.2), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide4, Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
             "💡 Key Takeaway", 18, True, AMBER)
add_text_box(slide4, Inches(0.7), Inches(5.3), Inches(11.8), Inches(1.2),
             "Random Forest with GridSearchCV tuning achieves R²=0.967, explaining 96.7% of cereal yield variance. "
             "The dominant predictor is Governorate (28%), which acts as a proxy for unmeasured factors like soil type, "
             "irrigation access, and farming traditions. Among weather variables, the engineered Heat Index outperforms "
             "raw temperature, confirming that combined heat-humidity stress is the key meteorological driver.",
             14, False, BLACK)

# ============================================================
# SLIDE 5: Performance Metrics Deep Dive
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide5, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Performance Metrics — Tuned Random Forest", 30, True, WHITE, PP_ALIGN.LEFT)

# Big metric cards
metrics = [
    ("R² Score", "0.9673", "96.7% of variance\nexplained", GREEN_MED),
    ("RMSE", "451 T", "Average prediction\nerror in Tonnes", BLUE),
    ("MAE", "299 T", "Median absolute\nerror", AMBER),
    ("CV R²", "0.96 ± 0.02", "5-fold cross-validation\nminimal overfitting", GREEN_DARK),
]
x_positions = [0.5, 3.7, 6.9, 10.1]
for i, (label, value, desc, color) in enumerate(metrics):
    x = x_positions[i]
    add_shape_bg(slide5, Inches(x), Inches(1.5), Inches(2.8), Inches(2.5), GREEN_PALE)
    add_shape_bg(slide5, Inches(x), Inches(1.5), Inches(2.8), Inches(0.1), color)
    add_text_box(slide5, Inches(x), Inches(1.7), Inches(2.8), Inches(0.4),
                 label, 14, False, GRAY, PP_ALIGN.CENTER)
    add_text_box(slide5, Inches(x), Inches(2.2), Inches(2.8), Inches(0.8),
                 value, 36, True, color, PP_ALIGN.CENTER)
    add_text_box(slide5, Inches(x), Inches(3.2), Inches(2.8), Inches(0.7),
                 desc, 12, False, GRAY, PP_ALIGN.CENTER)

# Interpretation
add_text_box(slide5, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
             "📋 Interpretation", 20, True, GREEN_DARK)
interp_items = [
    "• R² = 0.967 → For every 100 Tonnes of real yield variation, our model captures 96.7 Tonnes",
    "• RMSE = 451 T → On average, predictions deviate by ~451 Tonnes from actual values",
    "• MAE = 299 T → Half of all predictions are within 299 Tonnes of the true yield",
    "• CV stability → The small standard deviation (±0.02) across folds confirms no overfitting",
    "• Train vs Test gap → R² Train (0.98) vs R² Test (0.97) shows excellent generalization"
]
add_bullet_list(slide5, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.5), interp_items, 14, BLACK)

# ============================================================
# SLIDE 6: Key Insights
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide6, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide6, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Key Insights & Findings", 30, True, WHITE, PP_ALIGN.LEFT)

insights = [
    ("🗺️ Region is King", "Governorate alone explains 28% of yield variance — it encodes soil quality, altitude, irrigation patterns, and traditional farming practices that weather alone cannot capture.", Inches(0.5), Inches(1.5)),
    ("🌡️ Heat > Rain", "The engineered Heat Index (temp × humidity factor) outperforms precipitation as a predictor, confirming that Tunisian cereal farming is more limited by heat stress than water scarcity.", Inches(6.8), Inches(1.5)),
    ("☀️ Solar Drives Growth", "Solar radiation ranks #3 among all features, reflecting its direct role in photosynthesis and crop development during the growing season.", Inches(0.5), Inches(3.8)),
    ("❌ Linearity Fails", "Linear Regression (R²=0.25) vs Random Forest (R²=0.97) proves that crop-climate relationships involve thresholds, interactions, and non-linear responses that only ensemble methods can capture.", Inches(6.8), Inches(3.8)),
]

for title, desc, x, y in insights:
    add_shape_bg(slide6, x, y, Inches(6), Inches(2.0), GREEN_PALE)
    add_text_box(slide6, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.5),
                 title, 18, True, GREEN_DARK)
    add_text_box(slide6, x + Inches(0.2), y + Inches(0.7), Inches(5.6), Inches(1.2),
                 desc, 13, False, BLACK)

# ============================================================
# SLIDE 7: Applications
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide7, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide7, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Real-World Applications", 30, True, WHITE, PP_ALIGN.LEFT)

apps = [
    ("👨‍🌾 For Farmers", [
        "Predict expected yield before harvest",
        "Optimize fertilizer & water allocation",
        "Compare performance across seasons",
        "Plan storage and logistics"
    ], Inches(0.5)),
    ("🏛️ For Policymakers", [
        "National food security forecasting",
        "Targeted subsidy allocation",
        "Climate adaptation planning",
        "Import/export decision support"
    ], Inches(4.6)),
    ("🔬 For Researchers", [
        "Identify climate-yield relationships",
        "Quantify feature importance",
        "Benchmark for advanced models",
        "Reproducible ML pipeline"
    ], Inches(8.7)),
]

for title, items, x in apps:
    add_shape_bg(slide7, x, Inches(1.5), Inches(3.8), Inches(5.3), GREEN_PALE)
    add_text_box(slide7, x + Inches(0.2), Inches(1.6), Inches(3.4), Inches(0.5),
                 title, 20, True, GREEN_DARK)
    add_bullet_list(slide7, x + Inches(0.2), Inches(2.3), Inches(3.4), Inches(4.0),
                    [f"• {item}" for item in items], 14, BLACK, Pt(10))

# ============================================================
# SLIDE 8: Conclusion & Next Steps
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide8, Inches(0), Inches(0), Inches(13.333), Inches(1.2), GREEN_DARK)
add_text_box(slide8, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "Conclusion & Next Steps", 30, True, WHITE, PP_ALIGN.LEFT)

# Conclusion
add_shape_bg(slide8, Inches(0.5), Inches(1.5), Inches(6), Inches(3.0), GREEN_PALE)
add_text_box(slide8, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.5),
             "✅ What We Achieved", 20, True, GREEN_DARK)
conclusion_items = [
    "• Comprehensive EDA on 1,479 observations",
    "• 4 engineered features from raw weather",
    "• 4 models trained & rigorously compared",
    "• Best model: Tuned RF (R² = 0.967)",
    "• Identified Governorate + Heat Index as top drivers",
    "• Built interactive Streamlit prediction app"
]
add_bullet_list(slide8, Inches(0.7), Inches(2.2), Inches(5.5), Inches(2.0), conclusion_items, 14, BLACK, Pt(4))

# Next steps
add_shape_bg(slide8, Inches(6.8), Inches(1.5), Inches(6), Inches(3.0), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide8, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
             "🔮 Next Steps", 20, True, BLUE)
next_items = [
    "• Add monthly weather resolution",
    "• Include soil type & irrigation data",
    "• Integrate NDVI satellite imagery",
    "• Extend to multi-target prediction",
    "• Explore LSTM time-series models",
    "• Deploy as cloud API service"
]
add_bullet_list(slide8, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.0), next_items, 14, BLACK, Pt(4))

# Thank you
add_shape_bg(slide8, Inches(2), Inches(4.9), Inches(9.3), Inches(2.0), GREEN_DARK)
add_text_box(slide8, Inches(2), Inches(5.0), Inches(9.3), Inches(0.8),
             "Thank You — Questions?", 36, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide8, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
             "Ferdaws Saidi  &  Aya Gharsalli  |  EcoCrop Tunisia  |  2024", 16, False, RGBColor(0xA5, 0xD6, 0xA7), PP_ALIGN.CENTER)
add_text_box(slide8, Inches(2), Inches(6.3), Inches(9.3), Inches(0.5),
             "🌾 Predicting the Future of Tunisian Agriculture", 14, True, GREEN_LIGHT, PP_ALIGN.CENTER)

# Save
output_path = "PRESENTATION.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Size: {os.path.getsize(output_path) / 1024:.0f} KB")
print(f"   Slides: 8")